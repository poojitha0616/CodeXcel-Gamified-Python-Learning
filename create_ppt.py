from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Bullet layout
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        
        tf = body_shape.text_frame
        tf.text = content_text_list[0]
        
        for item in content_text_list[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    # Helper for Title Slide
    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = subtitle_text

    # --- Slide 1: Title ---
    add_title_slide("CodeXcel", "Gamified Python Learning Environment\nLevel Up Your Coding Skills")

    # --- Slide 2: Problem Statement ---
    add_slide("Problem Statement", [
        "The Challenge:",
        "- Learning to code can be dry and intimidating.",
        "- Traditional tutorials lack immediate feedback.",
        "- High dropout rates due to low engagement.",
        "",
        "The Need:",
        "- An interactive platform that makes learning fun.",
        "- Instant gratification through rewards.",
        "- Competitive elements to drive improvement."
    ])

    # --- Slide 3: Solution ---
    add_slide("Solution: CodeXcel", [
        "What is CodeXcel?",
        "- A comprehensive web-based Python learning platform.",
        "- Combines coding challenges with game mechanics.",
        "- Provides a 'Play to Learn' experience.",
        "",
        "Core Value Proposition:",
        "- Learn: Structured modules and hands-on coding.",
        "- Compete: Global leaderboards and rankings.",
        "- Achieve: Earn XP, badges, and unlock themes."
    ])

    # --- Slide 4: Gamification Features ---
    add_slide("Key Features - Gamification", [
        "Engaging the User:",
        "- XP System: +50 XP per successful challenge.",
        "- Leveling Up: Progressive difficulty tiers.",
        "- Badges: Visual achievements (e.g., 'Problem Solver').",
        "- Visual Feedback: Particle bursts, confetti, animations.",
        "",
        "(Visuals: Dashboard XP bar, Level Up modal)"
    ])

    # --- Slide 5: Learning Environment ---
    add_slide("Key Features - Learning Environment", [
        "Interactive Coding:",
        "- Real-time Code Editor: Python syntax highlighting.",
        "- Instant Feedback: Run code and see output immediately.",
        "- Judge0 Integration: Robust, safe code execution.",
        "- Structured Modules: Basics to Advanced topics.",
        "",
        "(Visuals: Challenge Interface, Code Editor)"
    ])

    # --- Slide 6: Customization & Social ---
    add_slide("Key Features - Customization & Social", [
        "Personalization:",
        "- Themes: Unlockable UI themes (Dark Neon, Cyber Grid).",
        "- Avatars: Unlockable avatar styles.",
        "",
        "Social Competition:",
        "- Leaderboard: Real-time ranking based on XP.",
        "- Profile: Showcase stats and earned badges.",
        "",
        "(Visuals: Customization page, Leaderboard)"
    ])

    # --- Slide 7: Technical Architecture ---
    add_slide("Technical Architecture", [
        "The Stack:",
        "- Frontend: HTML5, CSS3, Bootstrap 5, Jinja2.",
        "- Backend: Python Flask (v2.3.3).",
        "- Database: MongoDB (NoSQL).",
        "- Execution: Judge0 API (via RapidAPI).",
        "",
        "Data Flow:",
        "User Code -> Flask -> Judge0 -> Flask -> MongoDB -> UI Update"
    ])

    # --- Slide 8: Database Schema ---
    add_slide("Database Schema (MongoDB)", [
        "Collections:",
        "- Users: Profile, auth, XP, level, badges.",
        "- Modules: Challenge groupings.",
        "- Challenges: Problem, input/output, rewards.",
        "- Submissions: User attempt history.",
        "",
        "Why MongoDB?",
        "- Flexible schema for evolving gamification features."
    ])

    # --- Slide 9: Admin Capabilities ---
    add_slide("Admin Capabilities", [
        "Management Portal:",
        "- Dashboard: Overview of system stats.",
        "- Content Management: CRUD for coding challenges.",
        "- User Oversight: Monitor progress.",
        "",
        "(Visuals: Admin Panel)"
    ])

    # --- Slide 10: Future Roadmap ---
    add_slide("Future Roadmap", [
        "What's Next?",
        "- Multi-language Support (JS, C++, Java).",
        "- Social Features (Comments, Friends).",
        "- Advanced Gamification (Streaks, Boss Battles).",
        "- Mobile App Development."
    ])

    # --- Slide 11: Conclusion ---
    add_slide("Conclusion", [
        "Summary:",
        "- Bridges the gap between education and entertainment.",
        "- Robust, scalable, modern platform.",
        "- Ready for deployment.",
        "",
        "Thank You!",
        "Questions?"
    ])

    prs.save('CodeXcel_Presentation.pptx')
    print("Presentation saved as CodeXcel_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
